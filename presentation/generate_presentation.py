#!/usr/bin/env python3
"""Generate Resend Invoice Demo presentation (resend-invoice-demo.pptx)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# ── Dimensions (16:9) ────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)
PAD = Inches(0.75)

# ── Brand palette ─────────────────────────────────────────────────────────────
BG       = RGBColor(0x0C, 0x0C, 0x0E)   # Iron / black
SURFACE  = RGBColor(0x11, 0x11, 0x13)   # Surface / card
BORDER   = RGBColor(0x2E, 0x2F, 0x37)   # Border / subtle
BODY     = RGBColor(0xED, 0xEE, 0xF0)   # Eggshell body text
SECONDARY= RGBColor(0x8B, 0x8D, 0x98)   # Stone secondary text
ACCENT   = RGBColor(0xFF, 0xFF, 0xFF)   # Pure white callouts
CODE_BG  = RGBColor(0x18, 0x19, 0x1B)   # Code block bg

# ── Fonts (system-safe substitutes; swap in brand fonts if installed) ─────────
F_DISPLAY = "Georgia"          # → Domaine Display
F_SANS    = "Helvetica Neue"   # → ABC Favorit
F_BODY    = "Helvetica Neue"   # → Inter
F_CODE    = "Menlo"            # → JetBrains Mono / Inter Mono

TOTAL = 9

# ── Helpers ───────────────────────────────────────────────────────────────────

def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def txt(slide, text, left, top, width, height,
        font=F_BODY, size=14, color=BODY,
        bold=False, italic=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return box


def bullets(slide, items, left, top, width, height, size=16):
    """Render bullet list; lines starting with two spaces are sub-bullets."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        is_sub = item.startswith("  ")
        text = item.lstrip()
        sym = "→" if is_sub else "•"
        r = p.add_run()
        r.text = f"  {sym}  {text}"
        r.font.name = F_BODY
        r.font.size = Pt(size - 2 if is_sub else size)
        r.font.color.rgb = SECONDARY if is_sub else BODY
        p.space_before = Pt(3 if is_sub else 7)
    return box


def divider(slide, top):
    s = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        PAD, top, W - PAD * 2, Inches(0.012)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = BORDER
    s.line.fill.background()


def code_block(slide, code, left, top, width, height, size=9.5):
    # Background
    bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = CODE_BG
    bg.line.color.rgb = BORDER
    bg.line.width = Pt(0.75)

    # Text
    box = slide.shapes.add_textbox(
        left + Inches(0.22), top + Inches(0.18),
        width - Inches(0.44), height - Inches(0.36)
    )
    tf = box.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.name = F_CODE
        r.font.size = Pt(size)
        r.font.color.rgb = BODY


def footer(slide, num):
    txt(slide, "Resend",
        Inches(0.5), H - Inches(0.42), Inches(1.8), Inches(0.32),
        font=F_DISPLAY, size=9, color=SECONDARY, bold=True)
    txt(slide, f"{num} / {TOTAL}",
        W - Inches(1.6), H - Inches(0.42), Inches(1.4), Inches(0.32),
        size=9, color=SECONDARY, align=PP_ALIGN.RIGHT)


def section_label(slide, label):
    txt(slide, label,
        PAD, Inches(0.38), Inches(5), Inches(0.32),
        size=8.5, color=SECONDARY)


def slide_title(slide, title, top=Inches(0.72)):
    txt(slide, title,
        PAD, top, W - PAD * 2, Inches(0.92),
        font=F_DISPLAY, size=36, color=ACCENT, bold=True)
    divider(slide, top + Inches(0.88))


# ── Build presentation ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]  # truly blank layout

# ═══ SLIDE 1 — Title ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s)

# Top accent bar
bar = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, Inches(0.045))
bar.fill.solid(); bar.fill.fore_color.rgb = BORDER; bar.line.fill.background()

txt(s, "Resend Invoice Demo",
    PAD, H * 0.22, W - PAD * 2, Inches(1.7),
    font=F_DISPLAY, size=54, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

txt(s, "An agentic approach to a transactional email service",
    PAD, H * 0.52, W - PAD * 2, Inches(0.8),
    size=20, color=SECONDARY, align=PP_ALIGN.CENTER)

divider(s, H * 0.68)

txt(s, "Built by: [Your Name]  ·  February 2026",
    PAD, H * 0.73, W - PAD * 2, Inches(0.5),
    size=13, color=SECONDARY, align=PP_ALIGN.CENTER)

footer(s, 1)

# ═══ SLIDE 2 — What I Built ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s)
section_label(s, "01  OVERVIEW")
slide_title(s, "What I Built")

bullets(s, [
    "A Node.js + Express service that sends invoice emails via Resend",
    "Single  POST /invoice  endpoint orchestrates the full flow:",
    "  Validate input → Generate PDF → Send email with attachment → Optionally schedule a receipt",
    "Webhook listener at  POST /webhooks/resend  handles all Resend delivery events",
    "Covers four core Resend capabilities: transactional email, attachments, scheduling, webhooks",
], PAD, Inches(1.82), W - PAD * 2, Inches(4.8))

footer(s, 2)

# ═══ SLIDE 3 — Agentic Approach ═══════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s)
section_label(s, "02  PROCESS")
slide_title(s, "How I Built It — The Agentic Approach")

bullets(s, [
    "Used AI agents (via Cursor) to implement all 7 phases of the project",
    "Broke the project into 14 discrete agent tasks across phases 1–7",
    "Used 3 sub-agents with unique skills: resend-implementer, verifier, webhook-specialist",
    "Each task had a structured prompt, clear inputs, and defined outputs",
    "Agents operated with dependencies — Phase N only started after Phase N-1 completed",
    "Result: full working service built with some manual coding",
], PAD, Inches(1.82), W - PAD * 2, Inches(4.8))

footer(s, 3)

# ═══ SLIDE 4 — Architecture ═══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s)
section_label(s, "03  ARCHITECTURE")
slide_title(s, "Architecture", top=Inches(0.72))

code_block(s, """\
POST /invoice
  { lineItems, clientName, clientEmail, schedule_receipt, delay_minutes }
       │
       ├── Validate inputs (400 on failure)
       ├── generateInvoicePDF()   →  PDFKit → Buffer
       ├── sendInvoiceEmail()     →  Resend SDK (with Base64 PDF attachment)
       └── scheduleReceiptEmail() →  Resend scheduled_at (ISO 8601)

POST /webhooks/resend
  → express.raw()  (unparsed body required for Svix HMAC verification)
  → resend.webhooks.verify()
  → switch on event type  (sent, delivered, bounced, complained, ...)""",
    PAD, Inches(1.72), W - PAD * 2, Inches(4.3), size=10.5)

txt(s, "Three source files:  src/index.js  ·  src/invoice.js  ·  src/email.js",
    PAD, Inches(6.15), W - PAD * 2, Inches(0.38),
    size=11, color=SECONDARY)

footer(s, 4)

# ═══ SLIDE 5 — Transactional Email + Attachment ═══════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s)
section_label(s, "04  FEATURE DEEP-DIVE")
slide_title(s, "Transactional Email + PDF Attachment")

bullets(s, [
    "Invoice email sent via Resend Node SDK with HTML body and inline styles",
    "PDF generated in-memory with PDFKit — no temp files written to disk",
    "Attachment passed as Base64:  pdfBuffer.toString('base64')",
    "Resend attachment format:  [{ filename: 'invoice.pdf', content: base64String }]",
    "Unique invoice IDs generated per request:  INV-YYYYMMDD-XXXX",
], PAD, Inches(1.82), W - PAD * 2, Inches(4.8))

footer(s, 5)

# ═══ SLIDE 6 — Scheduled Email ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s)
section_label(s, "05  FEATURE DEEP-DIVE")
slide_title(s, "Scheduled Email")

bullets(s, [
    "Optional receipt email sent after a configurable delay (in minutes)",
    "Uses Resend's  scheduled_at  field — ISO 8601 datetime",
    "Calculated as:  new Date(Date.now() + delayMinutes * 60000).toISOString()",
    "Triggered by  schedule_receipt: true  in the request payload",
    "Receipt email ID returned in API response for tracking",
], PAD, Inches(1.82), W - PAD * 2, Inches(4.8))

footer(s, 6)

# ═══ SLIDE 7 — Webhooks ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s)
section_label(s, "06  FEATURE DEEP-DIVE")
slide_title(s, "Webhooks")

bullets(s, [
    "Endpoint:  POST /webhooks/resend",
    "Signature verified via  resend.webhooks.verify()  — Svix under the hood",
    "Critical detail:  express.raw()  scoped to webhook route only — unparsed body required for HMAC verification",
    "Handles 11 event types: email.sent, email.delivered, email.bounced, email.complained, and more",
    "Returns  401  on invalid signature,  200  on all valid events",
], PAD, Inches(1.82), W - PAD * 2, Inches(4.8))

footer(s, 7)

# ═══ SLIDE 8 — Sample Request & Response ══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s)
section_label(s, "07  DEMO")
slide_title(s, "Sample Request & Response", top=Inches(0.72))

COL = (W - PAD * 3) / 2

txt(s, "Request", PAD, Inches(1.72), COL, Inches(0.32), size=10, color=SECONDARY)
code_block(s, """\
curl -X POST http://localhost:3000/invoice \\
  -H "Content-Type: application/json" \\
  -d '{
    "lineItems": [
      { "description": "Consulting",
        "quantity": 5, "rate": 150 },
      { "description": "Design review",
        "quantity": 2, "rate": 200 }
    ],
    "clientName": "Acme Corp",
    "clientEmail": "billing@acme.com",
    "schedule_receipt": true,
    "delay_minutes": 1
  }'""",
    PAD, Inches(2.08), COL, Inches(4.1), size=9)

txt(s, "Response", PAD * 2 + COL, Inches(1.72), COL, Inches(0.32), size=10, color=SECONDARY)
code_block(s, """\
{
  "success": true,
  "invoiceId": "INV-20260220-4823",
  "invoice_total": 1150,
  "from": "Resend <onboarding@resend.dev>",
  "to": "billing@acme.com",
  "scheduledEmailId": "re_abc123..."
}""",
    PAD * 2 + COL, Inches(2.08), COL, Inches(2.6), size=10)

footer(s, 8)

# ═══ SLIDE 9 — Takeaways ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_bg(s)
section_label(s, "08  TAKEAWAYS")
slide_title(s, "Takeaways & What's Next")

bullets(s, [
    "Resend's SDK is clean and predictable — attachment, scheduling, and webhook verify are first-class",
    "The  scheduled_at  API is a standout feature — no job queue needed for simple delays",
    "Svix signature verification requires raw body discipline — easy to get wrong, easy to document",
    "Agentic development with structured tasks dramatically accelerated delivery",
], PAD, Inches(1.82), W - PAD * 2, Inches(2.9))

txt(s, "Possible Extensions",
    PAD, Inches(4.6), W - PAD * 2, Inches(0.42),
    size=12, color=ACCENT, bold=True)

bullets(s, [
    "Store invoice records in a database",
    "Retry on bounce event",
    "PDF templates with branding",
    "Resend Broadcasts for bulk invoicing",
], PAD, Inches(5.05), W - PAD * 2, Inches(1.8), size=14)

footer(s, 9)

# ── Save ──────────────────────────────────────────────────────────────────────
OUTPUT = "/Users/jvajda/Documents/Github_Personal/resend-demo/resend-invoice-demo.pptx"
prs.save(OUTPUT)
print(f"Saved → {OUTPUT}")
